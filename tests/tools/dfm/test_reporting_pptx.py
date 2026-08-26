from pathlib import Path

import pytest


pptx = pytest.importorskip("pptx")


def _image(path: Path, color: tuple[int, int, int]) -> None:
    from PIL import Image

    Image.new("RGB", (640, 420), color).save(path)


def test_pptx_report_contains_bounded_ordered_evidence(tmp_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from tools.dfm.reporting import render_default_reports

    for name, color in (
        ("issue_dfm-001_front.png", (225, 235, 245)),
        ("issue_dfm-001_oblique.png", (215, 230, 220)),
        ("issue_dfm-001_section.png", (245, 225, 220)),
        ("issue_dfm-001_extra.png", (235, 225, 245)),
        ("overview.png", (240, 240, 240)),
    ):
        _image(tmp_path / name, color)

    result = {
        "input": str(tmp_path / "part.step"),
        "issue_count": 1,
        "stats": {"volume_mm3": 1250.5, "face_count": 24},
        "issues": [
            {
                "id": "DFM-001",
                "code": "WALL_THIN",
                "title": "局部壁厚不足",
                "severity": "high",
                "message": "检测到局部壁厚低于当前阈值。",
                "metric": {"measured_mm": 0.8, "threshold_mm": 1.2},
                "images": [
                    "issue_dfm-001_front.png",
                    "issue_dfm-001_oblique.png",
                    "issue_dfm-001_section.png",
                    "issue_dfm-001_extra.png",
                ],
            }
        ],
    }

    artifacts = render_default_reports(
        artifact_dir=tmp_path,
        result=result,
        process="injection",
        scope_id="injection.wall-draft",
    )

    assert len(artifacts) == 1
    assert artifacts[0].kind == "report_presentation"
    assert artifacts[0].path.name == "dfm_report.pptx"
    deck = Presentation(artifacts[0].path)
    assert len(deck.slides) == 5

    detail = deck.slides[3]
    detail_text = "\n".join(
        shape.text for shape in detail.shapes if hasattr(shape, "text_frame")
    )
    assert detail_text.index("正视") < detail_text.index("剖视") < detail_text.index("斜视")
    assert sum(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in detail.shapes) == 3
