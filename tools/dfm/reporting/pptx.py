"""Editable PowerPoint report renderer backed by python-pptx."""

from __future__ import annotations

from collections import Counter
import importlib.util
from pathlib import Path
from typing import Any, Iterable, Mapping

from .base import ReportArtifact, ReportContext, ReportRenderer
from ..errors import DFMError


PPTX_MEDIA_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)

_NAVY = "14213D"
_TEAL = "087E8B"
_INK = "17202A"
_MUTED = "667085"
_PAPER = "F7F8FA"
_WHITE = "FFFFFF"
_BORDER = "D9DEE7"
_SEVERITY = {
    "critical": "B42318",
    "high": "D92D20",
    "medium": "DC6803",
    "low": "1570A6",
    "info": "475467",
}
_VIEW_ORDER = (
    ("front", "正视"),
    ("section", "剖视"),
    ("oblique", "斜视"),
)


def pptx_available() -> bool:
    try:
        return importlib.util.find_spec("pptx") is not None
    except (ImportError, AttributeError, ValueError):
        return False


def _api() -> dict[str, Any]:
    try:
        from PIL import Image
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise DFMError(
            "dependency_missing",
            "python-pptx is required to generate the DFM PowerPoint report.",
            {"dependency": "python-pptx", "install_extra": "hermes-agent[dfm]"},
        ) from exc
    return {
        "Image": Image,
        "Presentation": Presentation,
        "RGBColor": RGBColor,
        "MSO_SHAPE": MSO_SHAPE,
        "MSO_ANCHOR": MSO_ANCHOR,
        "PP_ALIGN": PP_ALIGN,
        "Inches": Inches,
        "Pt": Pt,
    }


def _rgb(value: str):
    api = _api()
    return api["RGBColor"](*(int(value[index:index + 2], 16) for index in (0, 2, 4)))


def _text(
    slide,
    value: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 16,
    color: str = _INK,
    bold: bool = False,
    font: str = "Microsoft YaHei",
    align: str = "left",
    valign: str = "middle",
):
    api = _api()
    box = slide.shapes.add_textbox(
        api["Inches"](x), api["Inches"](y), api["Inches"](w), api["Inches"](h)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = api["Inches"](0.04)
    frame.margin_top = frame.margin_bottom = api["Inches"](0.02)
    frame.vertical_anchor = {
        "top": api["MSO_ANCHOR"].TOP,
        "bottom": api["MSO_ANCHOR"].BOTTOM,
    }.get(valign, api["MSO_ANCHOR"].MIDDLE)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = {
        "center": api["PP_ALIGN"].CENTER,
        "right": api["PP_ALIGN"].RIGHT,
    }.get(align, api["PP_ALIGN"].LEFT)
    run = paragraph.add_run()
    run.text = str(value)
    run.font.name = font
    run.font.size = api["Pt"](size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _shape(slide, x: float, y: float, w: float, h: float, fill: str, *, radius=True):
    api = _api()
    shape_type = api["MSO_SHAPE"].ROUNDED_RECTANGLE if radius else api["MSO_SHAPE"].RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        api["Inches"](x), api["Inches"](y), api["Inches"](w), api["Inches"](h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.fill.background()
    return shape


def _base(slide, *, dark: bool = False) -> None:
    _shape(slide, 0, 0, 13.333, 7.5, _NAVY if dark else _PAPER, radius=False)


def _title(slide, title: str, subtitle: str | None = None) -> None:
    _text(slide, title, 0.65, 0.32, 10.8, 0.55, size=26, bold=True, valign="top")
    if subtitle:
        _text(slide, subtitle, 0.67, 0.9, 11.7, 0.3, size=10, color=_MUTED)


def _footer(slide, page: int, scope_id: str) -> None:
    _text(slide, f"Hermes DFM · {scope_id}", 0.65, 7.14, 5.5, 0.2, size=8, color=_MUTED)
    _text(slide, str(page), 12.05, 7.14, 0.55, 0.2, size=8, color=_MUTED, align="right")


def _safe_image(artifact_dir: Path, raw: Any) -> Path | None:
    if not isinstance(raw, str) or not raw.strip():
        return None
    relative = Path(raw)
    if relative.is_absolute():
        return None
    candidate = (artifact_dir / relative).resolve()
    if (
        not candidate.is_relative_to(artifact_dir.resolve())
        or not candidate.is_file()
        or candidate.suffix.lower() not in {".png", ".jpg", ".jpeg"}
    ):
        return None
    return candidate


def _issue_images(context: ReportContext, issue: Mapping[str, Any]) -> list[tuple[str, Path]]:
    raw_images = issue.get("images")
    values = list(raw_images) if isinstance(raw_images, list) else []
    if issue.get("image") and issue.get("image") not in values:
        values.append(issue.get("image"))
    resolved = [(str(raw), _safe_image(context.artifact_dir, raw)) for raw in values]
    ordered: list[tuple[str, Path]] = []
    used: set[Path] = set()
    for token, label in _VIEW_ORDER:
        match = next(
            (
                path
                for raw, path in resolved
                if path is not None and path not in used and token in Path(raw).stem.lower()
            ),
            None,
        )
        if match is not None:
            used.add(match)
            ordered.append((label, match))
    for _raw, path in resolved:
        if path is not None and path not in used:
            ordered.append(("证据", path))
            used.add(path)
    return ordered[: context.max_images_per_finding]


def _picture_contain(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    api = _api()
    with api["Image"].open(path) as image:
        image_w, image_h = image.size
    if image_w <= 0 or image_h <= 0:
        return
    scale = min(w / image_w, h / image_h)
    draw_w, draw_h = image_w * scale, image_h * scale
    slide.shapes.add_picture(
        str(path),
        api["Inches"](x + (w - draw_w) / 2),
        api["Inches"](y + (h - draw_h) / 2),
        width=api["Inches"](draw_w),
        height=api["Inches"](draw_h),
    )


def _short_value(value: Any, limit: int = 80) -> str:
    if isinstance(value, float):
        rendered = f"{value:.4g}"
    elif isinstance(value, (str, int, bool)) or value is None:
        rendered = str(value)
    elif isinstance(value, (list, tuple)) and len(value) <= 4:
        rendered = ", ".join(_short_value(item, 20) for item in value)
    else:
        rendered = ""
    return rendered if len(rendered) <= limit else rendered[: limit - 1] + "…"


def _metric_rows(issue: Mapping[str, Any]) -> list[tuple[str, str]]:
    metric = issue.get("metric")
    if not isinstance(metric, Mapping):
        return []
    rows = []
    for key, value in metric.items():
        rendered = _short_value(value)
        if rendered:
            rows.append((str(key), rendered))
        if len(rows) == 6:
            break
    return rows


def _valid_issues(result: Mapping[str, Any]) -> list[Mapping[str, Any]]:
    issues = result.get("issues")
    if not isinstance(issues, list):
        return []
    return [item for item in issues if isinstance(item, Mapping)]


class PythonPptxReportRenderer(ReportRenderer):
    key = "pptx"

    def render(self, context: ReportContext) -> ReportArtifact:
        api = _api()
        presentation = api["Presentation"]()
        presentation.slide_width = api["Inches"](13.333)
        presentation.slide_height = api["Inches"](7.5)
        layout = presentation.slide_layouts[6]
        issues = _valid_issues(context.result)
        evidence = [(issue, _issue_images(context, issue)) for issue in issues]
        evidence = [(issue, images) for issue, images in evidence if images]

        self._cover(presentation.slides.add_slide(layout), context)
        self._summary(presentation.slides.add_slide(layout), context, issues)
        self._overview(presentation.slides.add_slide(layout), context, issues, evidence)
        for issue, images in evidence:
            self._finding(presentation.slides.add_slide(layout), context, issue, images)
        self._notes(presentation.slides.add_slide(layout), context, issues, evidence)
        for page, slide in enumerate(presentation.slides, start=1):
            if page != 1:
                _footer(slide, page, context.scope_id)

        output = context.artifact_dir / "dfm_report.pptx"
        try:
            presentation.save(output)
            reopened = api["Presentation"](output)
        except (OSError, ValueError) as exc:
            raise DFMError(
                "report_generation_failed",
                "The DFM PowerPoint report could not be written.",
                {"path": str(output)},
            ) from exc
        expected = 4 + len(evidence)
        if not output.is_file() or output.stat().st_size == 0 or len(reopened.slides) != expected:
            raise DFMError(
                "report_generation_failed",
                "The DFM PowerPoint report failed structural validation.",
                {"expected_slides": expected, "actual_slides": len(reopened.slides)},
            )
        return ReportArtifact("report_presentation", output, PPTX_MEDIA_TYPE)

    @staticmethod
    def _cover(slide, context: ReportContext) -> None:
        _base(slide, dark=True)
        _shape(slide, 0.68, 0.62, 0.16, 5.98, _TEAL, radius=False)
        _text(slide, "DFM 分析报告", 1.15, 1.2, 5.2, 0.85, size=38, color=_WHITE, bold=True)
        _text(slide, "注塑成型 · 可制造性评估", 1.18, 2.12, 4.9, 0.42, size=17, color="B8DDE1")
        input_name = Path(str(context.result.get("input") or "未命名零件")).name
        _text(slide, input_name, 1.18, 3.0, 5.2, 0.7, size=21, color=_WHITE, bold=True, valign="top")
        _text(slide, f"工艺：{context.process}\n规则范围：{context.scope_id}", 1.18, 4.0, 5.2, 0.9, size=12, color="CDD5DF", valign="top")
        image = _safe_image(context.artifact_dir, "overview.png") or _safe_image(context.artifact_dir, "model.png")
        _shape(slide, 7.0, 0.72, 5.55, 5.95, _WHITE)
        if image:
            _picture_contain(slide, image, 7.2, 0.95, 5.15, 5.5)

    @staticmethod
    def _summary(slide, context: ReportContext, issues: list[Mapping[str, Any]]) -> None:
        _base(slide)
        _title(slide, "分析摘要", "结果来自几何后端的确定性计算及 Hermes 生成的证据")
        counts = Counter(str(item.get("severity") or "info").lower() for item in issues)
        cards = [
            ("问题总数", len(issues), _NAVY),
            ("高风险", counts["critical"] + counts["high"], _SEVERITY["high"]),
            ("中风险", counts["medium"], _SEVERITY["medium"]),
            ("低风险", counts["low"], _SEVERITY["low"]),
        ]
        for index, (label, value, color) in enumerate(cards):
            x = 0.7 + index * 3.08
            _shape(slide, x, 1.4, 2.75, 1.35, _WHITE)
            _text(slide, str(value), x + 0.18, 1.55, 1.2, 0.65, size=31, color=color, bold=True)
            _text(slide, label, x + 0.2, 2.2, 2.15, 0.3, size=11, color=_MUTED)
        overview = _safe_image(context.artifact_dir, "overview.png") or _safe_image(context.artifact_dir, "model.png")
        _shape(slide, 0.7, 3.05, 6.15, 3.6, _WHITE)
        if overview:
            _picture_contain(slide, overview, 0.9, 3.25, 5.75, 3.2)
        _shape(slide, 7.15, 3.05, 5.48, 3.6, _WHITE)
        stats = context.result.get("stats")
        rows: list[tuple[str, str]] = []
        if isinstance(stats, Mapping):
            for key, value in stats.items():
                rendered = _short_value(value)
                if rendered:
                    rows.append((str(key), rendered))
                if len(rows) == 8:
                    break
        _text(slide, "模型信息", 7.45, 3.3, 2.2, 0.35, size=17, bold=True)
        if not rows:
            _text(slide, "报告中未提供可直接展示的模型统计字段。", 7.45, 3.9, 4.7, 0.5, size=12, color=_MUTED)
        for index, (key, value) in enumerate(rows):
            y = 3.9 + index * 0.31
            _text(slide, key, 7.45, y, 2.1, 0.24, size=9, color=_MUTED)
            _text(slide, value, 9.55, y, 2.6, 0.24, size=9, bold=True)

    @staticmethod
    def _overview(
        slide,
        context: ReportContext,
        issues: list[Mapping[str, Any]],
        evidence: list[tuple[Mapping[str, Any], list[tuple[str, Path]]]],
    ) -> None:
        _base(slide)
        _title(slide, "问题分布", f"{len(evidence)} 个问题包含可视化证据详情页")
        codes = Counter(str(item.get("code") or item.get("title") or "未分类") for item in issues)
        top = codes.most_common(8)
        _shape(slide, 0.7, 1.35, 7.6, 5.45, _WHITE)
        _text(slide, "主要问题类型", 1.0, 1.65, 2.5, 0.4, size=18, bold=True)
        max_count = max((count for _, count in top), default=1)
        for index, (code, count) in enumerate(top):
            y = 2.25 + index * 0.51
            _text(slide, code, 1.0, y, 2.45, 0.25, size=10)
            _shape(slide, 3.55, y + 0.03, max(0.15, 3.65 * count / max_count), 0.18, _TEAL, radius=False)
            _text(slide, str(count), 7.35, y, 0.45, 0.25, size=10, bold=True, align="right")
        _shape(slide, 8.62, 1.35, 4.0, 5.45, _NAVY)
        _text(slide, "证据页规则", 8.98, 1.75, 2.8, 0.45, size=19, color=_WHITE, bold=True)
        notes = [
            "每个问题最多 3 张证据图",
            "固定顺序：正视、剖视、斜视",
            "检测指标直接来自分析 JSON",
            "未生成证据的问题保留在 JSON/MD",
        ]
        for index, note in enumerate(notes, start=1):
            _shape(slide, 9.0, 2.55 + (index - 1) * 0.82, 0.4, 0.4, _TEAL)
            _text(slide, str(index), 9.0, 2.55 + (index - 1) * 0.82, 0.4, 0.4, size=11, color=_WHITE, bold=True, align="center")
            _text(slide, note, 9.58, 2.48 + (index - 1) * 0.82, 2.55, 0.56, size=11, color="E7ECF3", valign="top")

    @staticmethod
    def _finding(
        slide,
        context: ReportContext,
        issue: Mapping[str, Any],
        images: list[tuple[str, Path]],
    ) -> None:
        _base(slide)
        severity = str(issue.get("severity") or "info").lower()
        color = _SEVERITY.get(severity, _SEVERITY["info"])
        issue_id = str(issue.get("id") or "DFM")
        title = str(issue.get("title") or issue.get("code") or "DFM 问题")
        _title(slide, f"{issue_id} · {title}", str(issue.get("code") or ""))
        _shape(slide, 11.45, 0.34, 1.15, 0.45, color)
        _text(slide, severity.upper(), 11.45, 0.34, 1.15, 0.45, size=10, color=_WHITE, bold=True, align="center")
        _shape(slide, 0.7, 1.35, 3.28, 5.45, _WHITE)
        _text(slide, "问题说明", 1.0, 1.68, 1.6, 0.35, size=16, bold=True)
        message = str(issue.get("message") or "未提供问题说明。")
        _text(slide, message, 1.0, 2.12, 2.66, 1.15, size=11, color=_INK, valign="top")
        _text(slide, "检测指标", 1.0, 3.45, 1.6, 0.35, size=16, bold=True)
        rows = _metric_rows(issue)
        if not rows:
            _text(slide, "无结构化指标", 1.0, 3.9, 2.5, 0.3, size=10, color=_MUTED)
        for index, (key, value) in enumerate(rows):
            y = 3.9 + index * 0.39
            _text(slide, key, 1.0, y, 1.28, 0.28, size=8.5, color=_MUTED)
            _text(slide, value, 2.25, y, 1.42, 0.28, size=8.5, bold=True, align="right")

        evidence_x, evidence_w = 4.25, 8.37
        gap = 0.18
        cell_w = (evidence_w - gap * (len(images) - 1)) / max(1, len(images))
        for index, (label, path) in enumerate(images):
            x = evidence_x + index * (cell_w + gap)
            _shape(slide, x, 1.35, cell_w, 5.45, _WHITE)
            _text(slide, label, x + 0.15, 1.58, cell_w - 0.3, 0.32, size=12, bold=True, align="center")
            _picture_contain(slide, path, x + 0.12, 2.0, cell_w - 0.24, 4.5)

    @staticmethod
    def _notes(
        slide,
        context: ReportContext,
        issues: list[Mapping[str, Any]],
        evidence: Iterable[tuple[Mapping[str, Any], list[tuple[str, Path]]]],
    ) -> None:
        _base(slide, dark=True)
        evidence_count = sum(1 for _ in evidence)
        _text(slide, "报告说明", 0.8, 0.72, 4.4, 0.7, size=32, color=_WHITE, bold=True)
        _shape(slide, 0.82, 1.68, 11.7, 4.55, "1D2D4B")
        notes = [
            f"本报告覆盖 {len(issues)} 个检测问题，其中 {evidence_count} 个包含证据详情页。",
            "PPT 中的工程数值、等级和说明直接来自 dfm_report.json，报告层不新增推断。",
            "未展示证据图的问题仍完整保存在 JSON 和 Markdown artifact 中。",
            "分析结论用于首轮可制造性筛查；最终设计决策仍需结合材料、模具方案和企业标准复核。",
        ]
        for index, note in enumerate(notes, start=1):
            y = 2.08 + (index - 1) * 0.86
            _text(slide, f"0{index}", 1.25, y, 0.65, 0.42, size=16, color="5ED0D9", bold=True)
            _text(slide, note, 2.15, y - 0.04, 9.55, 0.58, size=13, color="E7ECF3", valign="top")
